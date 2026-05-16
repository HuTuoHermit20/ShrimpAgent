import os
import time
import fitz
import torch
import numpy as np  # 新增：用于向量平均
from openai import OpenAI
from dotenv import load_dotenv
from typing import List, Dict
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from sentence_transformers import SentenceTransformer

load_dotenv()


class DEEPSEEK_V4_FLASH_LLMCLIENT:
    def __init__(self, model: str = None, apiKey: str = None,
                 baseUrl: str = None, timeout: int = None):
        self.model = model or os.getenv("LLM_MODEL_ID")
        apiKey = apiKey or os.getenv("LLM_API_KEY")
        baseUrl = baseUrl or os.getenv("LLM_BASE_URL")
        timeout = timeout or int(os.getenv("LLM_TIMEOUT", 60))

        if not all([self.model, apiKey, baseUrl]):
            raise ValueError("您未在.env文件中提供LLM的模型ID、apikey和baseurl")

        self.client = OpenAI(api_key=apiKey, base_url=baseUrl, timeout=timeout)

    def think(self, messages: List[Dict[str, str]],
              temperature: float = 0) -> str:
        print(f"正在调用 {self.model} 模型...")
        start_time = time.monotonic_ns()
        try:
            response = self.client.chat.completions.create(
                model=self.model, messages=messages,
                temperature=temperature, stream=True,
            )
            print("大语言模型响应成功:")
            end_time = time.monotonic_ns()
            part_time = end_time - start_time
            print(f"本次调用共花费 {part_time} 纳秒")
            collected_content = []
            for chunk in response:
                if not chunk.choices:
                    continue
                content = chunk.choices[0].delta.content or ""
                print(content, end="", flush=True)
                collected_content.append(content)
            print()
            return "".join(collected_content)
        except Exception as e:
            print(f"调用LLM API时发生错误: {e}")
            return None


class FileEmbedder:
    def __init__(self, model_name_or_path="Qwen3-Embedding-0-6B", device=None):
        # 单卡环境直接使用字符串，避免 torch.device 缺索引报错
        self.device = 'cuda' if torch.cuda.is_available() else 'cpu'
        print(f"加载模型到 {self.device} ...")

        # 必须在CUDA初始化前设置环境变量
        os.environ['PYTORCH_CUDA_ALLOC_CONF'] = 'expandable_segments:True'

        self.model = SentenceTransformer(model_name_or_path, device=self.device)
        self.model.half()  # FP16 节省约一半显存
        # ✅ 移除 max_memory_gb 硬性限制！8G显卡让PyTorch自动管理最稳定

    @staticmethod
    def _extract_text(file_path):
        ext = os.path.splitext(file_path)[1].lower()
        if ext == '.docx':
            doc = Document(file_path)
            text_parts = [p.text.strip() for p in doc.paragraphs if
                          p.text.strip()]
            for table in doc.tables:
                for row in table.rows:
                    row_cells = [cell.text.strip() for cell in row.cells if
                                 cell.text.strip()]
                    if row_cells:
                        text_parts.append(' | '.join(row_cells))
            return '\n'.join(text_parts)
        elif ext == '.pptx':
            prs = Presentation(file_path)
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)
            return '\n'.join(text_parts)
        elif ext == '.xlsx':
            wb = load_workbook(file_path, data_only=True)
            rows_text = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    row_str = ' | '.join(str(c) for c in row if c is not None)
                    if row_str.strip():
                        rows_text.append(row_str)
            return '\n'.join(rows_text)
        elif ext == '.pdf':
            doc = fitz.open(file_path)
            text_parts = [page.get_text().strip() for page in doc if
                          page.get_text().strip()]
            doc.close()
            return '\n'.join(text_parts)
        else:
            raise ValueError(f"不支持的文件类型: {ext}")

    @staticmethod
    def _chunk_text(text: str, max_chars: int = 1200) -> List[str]:
        """将长文本切分为安全长度，防止单条塞入模型导致 OOM"""
        if len(text) <= max_chars:
            return [text]
        return [text[i:i + max_chars] for i in range(0, len(text), max_chars)]

    def embed_file(self, file_path, as_query=False):
        text = self._extract_text(file_path)
        if not text.strip():
            print(f"警告: '{file_path}' 无文本，跳过。")
            return text, None

        chunks = self._chunk_text(text)
        torch.cuda.empty_cache()  # 编码前释放缓存

        # batch_size=1 配合切分，稳过 8G 显存
        embeddings = self.model.encode(
            chunks,
            batch_size=1,
            show_progress_bar=False,
            prompt_name="query" if as_query else None
        )
        # 多段向量求平均，保留完整文档语义
        avg_embedding = np.mean(embeddings, axis=0).tolist()
        return text, avg_embedding

    def embed_files(self, file_paths, as_query=False):
        results = []
        for fp in file_paths:
            try:
                text = self._extract_text(fp)
                if not text.strip():
                    print(f"警告: '{fp}' 无文本。")
                    results.append((fp, None, None))
                    continue

                chunks = self._chunk_text(text)
                torch.cuda.empty_cache()
                embeddings = self.model.encode(
                    chunks, batch_size=1, show_progress_bar=False,
                    prompt_name="query" if as_query else None
                )
                avg_embedding = np.mean(embeddings, axis=0).tolist()
                results.append((fp, text, avg_embedding))
            except Exception as e:
                print(f"处理 {fp} 时出错: {e}")
                results.append((fp, None, None))
        return results

    def embed_directory(self, dir_path, as_query=False):
        supported_ext = {'.docx', '.pptx', '.xlsx', '.pdf'}
        file_paths = []
        for root, _, files in os.walk(dir_path):
            for file in files:
                if os.path.splitext(file)[1].lower() in supported_ext:
                    file_paths.append(os.path.join(root, file))
        if not file_paths:
            print(f"警告: 目录 '{dir_path}' 中未找到支持的文件。")
            return []
        return self.embed_files(file_paths, as_query)
